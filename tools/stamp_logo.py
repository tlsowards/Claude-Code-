#!/usr/bin/env python3
"""Stamp a logo onto every slide of one or more PowerPoint decks.

Typical use — brand a folder of decks before publishing them to a website:

    python3 tools/stamp_logo.py --logo logo.png --in decks/ --out branded/

Re-running on an already-stamped deck is safe: logos added by this script are
named so they can be found and replaced rather than piling up.
"""

import argparse
import io
import pathlib
import sys

from PIL import Image
from pptx import Presentation
from pptx.util import Emu

# Marks the shapes this script owns, so a re-run replaces them instead of
# stacking a second logo on top of the first.
SHAPE_NAME = "BrandLogo"

CORNERS = ("top-left", "top-right", "bottom-left", "bottom-right")


def trim_transparent(img):
    """Crop fully transparent borders so the logo sits tight against the margin."""
    if img.mode != "RGBA":
        return img
    bbox = img.getchannel("A").getbbox()
    return img.crop(bbox) if bbox else img


def load_logo(path, trim, opacity):
    img = Image.open(path)
    img = img.convert("RGBA")
    if trim:
        img = trim_transparent(img)
    if opacity < 1.0:
        alpha = img.getchannel("A").point(lambda a: int(a * opacity))
        img.putalpha(alpha)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf, img.width / img.height


def place(slide_w, slide_h, logo_ratio, width_pct, margin_pct, position):
    width = int(slide_w * width_pct / 100)
    height = int(width / logo_ratio)
    # Both insets come off slide width so the gap looks equal on screen,
    # rather than a 16:9 slide showing a visibly tighter bottom margin.
    margin = int(slide_w * margin_pct / 100)

    left = margin if position.endswith("left") else slide_w - width - margin
    top = margin if position.startswith("top") else slide_h - height - margin
    return Emu(left), Emu(top), Emu(width), Emu(height)


def clear_previous(slide):
    for shape in list(slide.shapes):
        if shape.name == SHAPE_NAME:
            shape._element.getparent().remove(shape._element)


def send_behind_text(slide, pic):
    """Drop the logo just above the background so it never covers body text."""
    tree = pic._element.getparent()
    tree.remove(pic._element)
    # A spTree always opens with nvGrpSpPr then grpSpPr, so index 2 is the
    # first drawable slot — behind every real shape, but still in the tree.
    tree.insert(2, pic._element)


def stamp(deck_path, out_path, logo_buf, logo_ratio, args):
    prs = Presentation(str(deck_path))
    slide_w, slide_h = prs.slide_width, prs.slide_height
    left, top, width, height = place(
        slide_w, slide_h, logo_ratio, args.width_pct, args.margin_pct, args.position
    )

    stamped = 0
    for index, slide in enumerate(prs.slides):
        if args.skip_first and index == 0:
            continue
        clear_previous(slide)
        # add_picture consumes the stream, so rewind it for each slide.
        # python-pptx hashes the bytes, so the deck still stores one copy.
        logo_buf.seek(0)
        pic = slide.shapes.add_picture(logo_buf, left, top, width, height)
        pic.name = SHAPE_NAME
        if args.behind:
            send_behind_text(slide, pic)
        stamped += 1

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return stamped, len(prs.slides)


def main():
    p = argparse.ArgumentParser(description=__doc__,
                                formatter_class=argparse.RawDescriptionHelpFormatter)
    p.add_argument("--logo", required=True, type=pathlib.Path,
                   help="Logo image (PNG with transparency works best)")
    p.add_argument("--in", dest="source", required=True, type=pathlib.Path,
                   help="A .pptx file, or a folder of them")
    p.add_argument("--out", required=True, type=pathlib.Path,
                   help="Output folder for the branded copies")
    p.add_argument("--position", default="bottom-right", choices=CORNERS)
    p.add_argument("--width-pct", type=float, default=12.0,
                   help="Logo width as %% of slide width (default: 12)")
    p.add_argument("--margin-pct", type=float, default=3.0,
                   help="Margin from the slide edge as %% of slide width (default: 3)")
    p.add_argument("--opacity", type=float, default=1.0,
                   help="0.0-1.0; below 1 makes a subtle watermark (default: 1)")
    p.add_argument("--skip-first", action="store_true",
                   help="Leave the title slide untouched")
    p.add_argument("--trim", action="store_true",
                   help="Crop transparent padding around the logo first")
    p.add_argument("--behind", action="store_true",
                   help="Place the logo behind slide text instead of on top")
    args = p.parse_args()

    if not args.logo.is_file():
        sys.exit(f"Logo not found: {args.logo}")
    if not args.source.exists():
        sys.exit(f"Not found: {args.source}")
    if not 0.0 < args.opacity <= 1.0:
        sys.exit("--opacity must be greater than 0 and at most 1")
    if not 0 < args.width_pct <= 100:
        sys.exit("--width-pct must be greater than 0 and at most 100")
    if not 0 <= args.margin_pct < 50:
        sys.exit("--margin-pct must be at least 0 and less than 50")

    out_dir = args.out.resolve()
    if args.source.is_dir():
        source = args.source.resolve()
        if out_dir == source:
            sys.exit("--out must differ from --in, so the originals are kept")
        # --out may sit inside --in; skip anything already there so a second
        # run brands the originals again rather than the branded copies.
        decks = sorted(
            d for d in args.source.rglob("*")
            if d.suffix.lower() == ".pptx"
            and not d.name.startswith("~$")
            and out_dir not in d.resolve().parents
        )
    else:
        decks = [args.source]
    if not decks:
        sys.exit(f"No .pptx files found in {args.source}")

    logo_buf, logo_ratio = load_logo(args.logo, args.trim, args.opacity)

    failed = 0
    for deck in decks:
        rel = deck.relative_to(args.source) if args.source.is_dir() else deck.name
        out_path = args.out / rel
        try:
            stamped, total = stamp(deck, out_path, logo_buf, logo_ratio, args)
        except Exception as exc:
            # One unreadable deck shouldn't abandon the rest of the batch.
            failed += 1
            print(f"{deck.name}: SKIPPED ({exc})", file=sys.stderr)
            continue
        print(f"{deck.name}: logo on {stamped}/{total} slides -> {out_path}")

    if failed:
        print(f"\n{failed} of {len(decks)} deck(s) could not be processed.",
              file=sys.stderr)
        return 1
    return 0


if __name__ == "__main__":
    sys.exit(main())
